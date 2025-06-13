from PIL import Image
import pytesseract
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

# ✅ OCR pour les images
def ocr_from_image(image_path):
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)

# ✅ Extraction PDF avec PyMuPDF
def extract_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# ✅ Lecture Word DOCX
def extract_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

# ✅ Lecture Excel XLSX
def extract_from_excel(path):
    df = pd.read_excel(path)
    return df.to_string(index=False)

# ✅ Lecture PowerPoint PPTX
def extract_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text

# ✅ Routeur général selon l'extension du fichier
def extract_text_from_file(file_path, extension):
    extension = extension.lower()
    if extension == ".pdf":
        return extract_from_pdf(file_path)
    elif extension in [".jpg", ".jpeg", ".png"]:
        return ocr_from_image(file_path)
    elif extension == ".docx":
        return extract_from_docx(file_path)
    elif extension == ".xlsx":
        return extract_from_excel(file_path)
    elif extension == ".pptx":
        return extract_from_pptx(file_path)
    else:
        return "❌ Format non supporté."
